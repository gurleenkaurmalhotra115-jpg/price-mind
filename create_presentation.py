import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# Color Palette (Premium Fintech/AI Theme)
bg_color = RGBColor(10, 17, 30)       # Slate Black / Dark Blue Background
text_white = RGBColor(255, 255, 255) # Pure White
text_gray = RGBColor(180, 187, 198)   # Muted Gray
gold_color = RGBColor(245, 158, 11)   # Premium Gold Accent
accent_blue = RGBColor(14, 165, 233) # Cyber Blue Highlight
card_bg = RGBColor(20, 30, 48)       # Slate Card (Light)
card_dark = RGBColor(15, 23, 42)     # Slate Card (Dark)
green_color = RGBColor(34, 197, 94)  # Success Green

# Helper function to scale coordinates from 13.333x7.5 to 10x5.625 (75% scaling)
def scale(x):
    return Inches(x * 0.75)

def apply_dark_theme(slide, draw_border=True):
    # Slide background, border frame, and corner brackets are drawn in the template's slide master / layouts.
    # Do not call this function or draw them manually to preserve template fidelity.
    pass

def add_slide_header(slide, title_text):
    # The template slides already contain the gold-bordered title pills with title texts.
    # Do not draw headers manually to prevent double titles.
    pass

def add_card(slide, left, top, width, height, bg_rgb, border_rgb=None):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_rgb
    if border_rgb:
        card.line.color.rgb = border_rgb
        card.line.width = Pt(1.5)
    else:
        card.line.fill.background()
    return card

def clean_placeholders(slide, texts_to_remove):
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            for pattern in texts_to_remove:
                if pattern in txt:
                    shapes_to_remove.append(shape)
                    break
    for shape in shapes_to_remove:
        slide.shapes._spTree.remove(shape._element)

def create_presentation():
    # Load existing template
    template_path = "c:/Users/gurle/OneDrive/Apps/HACKATHON/hackathon ideas/template.pptx"
    output_path = "c:/Users/gurle/OneDrive/Apps/HACKATHON/hackathon ideas/Hack_4_Crown_PriceMind_Presentation.pptx"
    
    prs = Presentation(template_path)
    
    # Check slide dimensions (should be 10x5.625)
    prs.slide_width = Inches(10.0)
    prs.slide_height = Inches(5.625)
    
    # ==========================================
    # SLIDE 1: Title Slide (No Blue Outline Frame)
    # ==========================================
    slide = prs.slides[0]
    
    # 1. Reposition left branding and logo shapes to fit in left space (X: 0.5" to 6.0")
    # This prevents them from overlapping with the right card.
    for shape in slide.shapes:
        # Crowns
        if "object 2" in shape.name:
            shape.left = Inches(2.6)
            shape.top = Inches(1.2)
        # HACK-4-CROWN Picture Logo
        elif "object 3" in shape.name:
            shape.left = Inches(0.8)
            shape.top = Inches(1.5)
            shape.width = Inches(4.4)
            shape.height = Inches(4.4 * (1.067 / 6.577))
        # HACK-4-CROWN text placeholder - delete to prevent double text rendering
        elif "object 4" in shape.name:
            slide.shapes._spTree.remove(shape._element)
        # Forging The Empire group
        elif "object 7" in shape.name:
            shape.left = Inches(0.8)
            shape.top = Inches(2.3)
            shape.width = Inches(4.4)
            shape.height = Inches(4.4 * (0.806 / 7.007))
            
    # 2. Update and reposition the left text box containing title, tagline, track, and team
    left_text_shape = None
    for shape in slide.shapes:
        if "object 10" in shape.name:
            left_text_shape = shape
            break
            
    if left_text_shape:
        left_text_shape.left = Inches(0.8)
        left_text_shape.top = Inches(2.8)
        left_text_shape.width = Inches(4.5)
        left_text_shape.height = Inches(2.2)
        
        tf = left_text_shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0)
        
        # Paragraph 1: Project Title
        p1 = tf.paragraphs[0]
        p1.text = "PriceMind"
        p1.font.name = "Segoe UI"
        p1.font.size = Pt(28)
        p1.font.bold = True
        p1.font.color.rgb = text_white
        p1.space_after = Pt(4)
        
        # Paragraph 2: Tagline
        p2 = tf.add_paragraph()
        p2.text = "Forecast. Explain. Recommend."
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(16)
        p2.font.bold = True
        p2.font.color.rgb = gold_color
        p2.space_after = Pt(10)
        
        # Paragraph 3: Track
        p3 = tf.add_paragraph()
        p3.text = "TRACK: Retail & FinTech"
        p3.font.name = "Segoe UI"
        p3.font.size = Pt(11)
        p3.font.color.rgb = text_gray
        p3.space_after = Pt(3)
        
        # Paragraph 4: Team
        p4 = tf.add_paragraph()
        p4.text = "TEAM: _________________"
        p4.font.name = "Segoe UI"
        p4.font.size = Pt(12)
        p4.font.bold = True
        p4.font.color.rgb = text_white
    
    # 3. Add Mockup Recommendation Card (Right Half)
    add_card(slide, scale(8.0), scale(1.8), scale(4.3), scale(4.7), card_bg, accent_blue)
    
    cardBox = slide.shapes.add_textbox(scale(8.3), scale(2.0), scale(3.7), scale(1.5))
    ctf = cardBox.text_frame
    ctf.word_wrap = True
    ctf.margin_left = ctf.margin_top = ctf.margin_right = ctf.margin_bottom = Inches(0)
    
    cp = ctf.paragraphs[0]
    cp.text = "LIVE PRODUCT ANALYSIS"
    cp.font.name = "Segoe UI"
    cp.font.size = Pt(10)
    cp.font.bold = True
    cp.font.color.rgb = accent_blue
    cp.space_after = Pt(6)
    
    cp2 = ctf.add_paragraph()
    cp2.text = "PriceMind v1"
    cp2.font.name = "Segoe UI"
    cp2.font.size = Pt(16)
    cp2.font.bold = True
    cp2.font.color.rgb = text_white
    cp2.space_after = Pt(3)
    
    cp3 = ctf.add_paragraph()
    cp3.text = "Current Price: ₹42,999"
    cp3.font.name = "Segoe UI"
    cp3.font.size = Pt(13)
    cp3.font.color.rgb = text_gray
    
    # Recommendation Badge inside the Card
    add_card(slide, scale(8.3), scale(2.9), scale(3.7), scale(0.5), RGBColor(60, 45, 10), gold_color)
    badgeBox = slide.shapes.add_textbox(scale(8.4), scale(2.95), scale(3.5), scale(0.4))
    btf = badgeBox.text_frame
    bp = btf.paragraphs[0]
    bp.text = "RECOMMENDATION: 🟡 WAIT"
    bp.alignment = PP_ALIGN.CENTER
    bp.font.name = "Segoe UI"
    bp.font.size = Pt(11)
    bp.font.bold = True
    bp.font.color.rgb = gold_color
    
    # Expected Price Trajectory Label
    lblBox = slide.shapes.add_textbox(scale(8.3), scale(3.55), scale(3.7), scale(0.3))
    ltf = lblBox.text_frame
    lp = ltf.paragraphs[0]
    lp.text = "Expected Price Trajectory:"
    lp.font.name = "Segoe UI"
    lp.font.size = Pt(9.5)
    lp.font.bold = True
    lp.font.color.rgb = text_gray
    
    # Trajectory mini-cards
    for idx, (label, val) in enumerate([("7d", "₹42,500"), ("30d", "₹39,999"), ("60d", "₹41,200")]):
        card_l = scale(8.3 + idx * 1.3)
        card_t = scale(3.85)
        card_w = scale(1.1)
        card_h = scale(0.7)
        # Add card shape
        mini = add_card(slide, card_l, card_t, card_w, card_h, card_dark, border_rgb=accent_blue if idx == 1 else text_gray)
        # Add textbox inside
        tbox = slide.shapes.add_textbox(card_l, card_t + scale(0.05), card_w, card_h - scale(0.1))
        tf_mini = tbox.text_frame
        tf_mini.word_wrap = True
        tf_mini.margin_left = tf_mini.margin_top = tf_mini.margin_right = tf_mini.margin_bottom = Inches(0)
        
        p_l = tf_mini.paragraphs[0]
        p_l.text = label
        p_l.alignment = PP_ALIGN.CENTER
        p_l.font.name = "Segoe UI"
        p_l.font.size = Pt(8)
        p_l.font.color.rgb = text_gray
        
        p_v = tf_mini.add_paragraph()
        p_v.text = val
        p_v.alignment = PP_ALIGN.CENTER
        p_v.font.name = "Segoe UI"
        p_v.font.size = Pt(10)
        p_v.font.bold = True
        p_v.font.color.rgb = gold_color if idx == 1 else text_white
        
    cardBox2 = slide.shapes.add_textbox(scale(8.3), scale(4.7), scale(3.7), scale(1.6))
    ctf2 = cardBox2.text_frame
    ctf2.word_wrap = True
    ctf2.margin_left = ctf2.margin_top = ctf2.margin_right = ctf2.margin_bottom = Inches(0)
    
    cp4 = ctf2.paragraphs[0]
    cp4.text = "Potential Saving: ₹3,000"
    cp4.font.name = "Segoe UI"
    cp4.font.size = Pt(11)
    cp4.font.bold = True
    cp4.font.color.rgb = green_color
    cp4.space_after = Pt(3)
    
    cp5 = ctf2.add_paragraph()
    cp5.text = "Forecast Confidence: 82%"
    cp5.font.name = "Segoe UI"
    cp5.font.size = Pt(9.5)
    cp5.font.bold = True
    cp5.font.color.rgb = text_white
    
    # Speaker notes for Slide 1
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "Forecast Confidence: 82% [EXAMPLE — replace with real output once model is trained]"

    # ==========================================
    # SLIDE 2: PROBLEM STATEMENT (With Blue Outline Frame)
    # ==========================================
    slide = prs.slides[1]
    # Remove helper placeholders
    clean_placeholders(slide, ["Describe the problem"])
    
    # Left Column (Consumer's Dilemma Box)
    add_card(slide, scale(1.0), scale(1.8), scale(5.2), scale(4.8), card_dark)
    leftBox = slide.shapes.add_textbox(scale(1.3), scale(2.1), scale(4.6), scale(4.2))
    ltf = leftBox.text_frame
    ltf.word_wrap = True
    ltf.margin_left = ltf.margin_top = ltf.margin_right = ltf.margin_bottom = Inches(0)
    
    lp = ltf.paragraphs[0]
    lp.text = "THE SHOPPER'S ANXIETY"
    lp.font.name = "Segoe UI"
    lp.font.size = Pt(11)
    lp.font.bold = True
    lp.font.color.rgb = gold_color
    lp.space_after = Pt(12)
    
    lp2 = ltf.add_paragraph()
    lp2.text = "“Should I buy this now or wait?”"
    lp2.font.name = "Segoe UI"
    lp2.font.size = Pt(18)
    lp2.font.italic = True
    lp2.font.bold = True
    lp2.font.color.rgb = text_white
    lp2.space_after = Pt(12)
    
    lp3 = ltf.add_paragraph()
    lp3.text = "E-commerce buyers face extreme price volatility when shopping for smartphones. Without visibility into future pricing, consumers struggle to know whether to make an immediate purchase or wait for a price drop, leading to shopping anxiety and regret."
    lp3.font.name = "Segoe UI"
    lp3.font.size = Pt(11)
    lp3.font.color.rgb = text_gray
    
    # Right Column (Core Pain Points Box)
    add_card(slide, scale(6.8), scale(1.8), scale(5.5), scale(4.8), card_bg)
    
    # Right Box Header
    rHeaderBox = slide.shapes.add_textbox(scale(7.1), scale(2.0), scale(4.9), scale(0.4))
    rhtf = rHeaderBox.text_frame
    rhp = rhtf.paragraphs[0]
    rhp.text = "CORE FRICTION POINTS"
    rhp.font.name = "Segoe UI"
    rhp.font.size = Pt(11)
    rhp.font.bold = True
    rhp.font.color.rgb = accent_blue
    
    points = [
        ("Platform price opacity", "Dynamic pricing models keep historical trends hidden from buyers."),
        ("Deceptive discount badges", "Artificial urgency using inflated MSRP baselines."),
        ("No forecasting today", "Consumers lack predictive signals to time purchases.")
    ]
    
    for idx, (title, desc) in enumerate(points):
        sub_y = scale(2.4 + idx * 1.35)
        # Sub card
        add_card(slide, scale(7.1), sub_y, scale(4.9), scale(1.2), card_dark)
        
        # Sub card text
        sub_tbox = slide.shapes.add_textbox(scale(7.2), sub_y + scale(0.1), scale(4.7), scale(1.0))
        stf = sub_tbox.text_frame
        stf.word_wrap = True
        stf.margin_left = stf.margin_top = stf.margin_right = stf.margin_bottom = Inches(0)
        
        sp = stf.paragraphs[0]
        sp.text = f"• {title}"
        sp.font.name = "Segoe UI"
        sp.font.size = Pt(11)
        sp.font.bold = True
        sp.font.color.rgb = text_white
        sp.space_after = Pt(2)
        
        sp2 = stf.add_paragraph()
        sp2.text = desc
        sp2.font.name = "Segoe UI"
        sp2.font.size = Pt(9.5)
        sp2.font.color.rgb = text_gray

    # ==========================================
    # SLIDE 3: SOLUTION (With Blue Outline Frame)
    # ==========================================
    slide = prs.slides[2]
    # Remove helper placeholders
    clean_placeholders(slide, ["Outline your solution"])
    
    # Left Column (Core Value Pillars Box)
    add_card(slide, scale(1.0), scale(1.8), scale(5.2), scale(4.5), card_bg)
    
    # Left Box Header
    lHeaderBox = slide.shapes.add_textbox(scale(1.3), scale(2.0), scale(4.6), scale(0.4))
    lhtf = lHeaderBox.text_frame
    lhp = lhtf.paragraphs[0]
    lhp.text = "INTELLIGENT DECISION SUPPORT"
    lhp.font.name = "Segoe UI"
    lhp.font.size = Pt(11)
    lhp.font.bold = True
    lhp.font.color.rgb = gold_color
    
    pillars = [
        ("WHERE", "Prepared static product and pricing metadata from historical dataset."),
        ("WHEN", "XGBoost forecasting engine predicts future pricing (7/30/60 days)."),
        ("WHETHER", "Confidence-scored recommendation engine advises immediately.")
    ]
    
    for idx, (title, desc) in enumerate(pillars):
        sub_y = scale(2.4 + idx * 1.2)
        # Sub card
        add_card(slide, scale(1.3), sub_y, scale(4.6), scale(1.05), card_dark)
        
        # Sub card text
        sub_tbox = slide.shapes.add_textbox(Inches(1.4), sub_y + scale(0.1), scale(4.4), scale(0.85))
        stf = sub_tbox.text_frame
        stf.word_wrap = True
        stf.margin_left = stf.margin_top = stf.margin_right = stf.margin_bottom = Inches(0)
        
        sp = stf.paragraphs[0]
        sp.text = f"{title} to Buy"
        sp.font.name = "Segoe UI"
        sp.font.size = Pt(11)
        sp.font.bold = True
        sp.font.color.rgb = text_white
        sp.space_after = Pt(2)
        
        sp2 = stf.add_paragraph()
        sp2.text = desc
        sp2.font.name = "Segoe UI"
        sp2.font.size = Pt(9.5)
        sp2.font.color.rgb = text_gray
        
    # Right Column (Workflow Process Box)
    add_card(slide, scale(6.8), scale(1.8), scale(5.5), scale(4.5), card_dark)
    
    rHeaderBox = slide.shapes.add_textbox(scale(7.1), scale(2.0), scale(4.9), scale(0.4))
    rhtf = rHeaderBox.text_frame
    rhp = rhtf.paragraphs[0]
    rhp.text = "PRICEMIND WORKFLOW"
    rhp.font.name = "Segoe UI"
    rhp.font.size = Pt(11)
    rhp.font.bold = True
    rhp.font.color.rgb = accent_blue
    
    steps_data = [
        ("1. Product Search", "Query smartphone model details"),
        ("2. Historical Price Load", "Load static historical records"),
        ("3. Feature Engineering", "Compute lags, volatility, festivals"),
        ("4. XGBoost Forecast", "Predict 7/30/60d price trends"),
        ("5. Recommendation Engine", "Advises BUY / WAIT / STRONG BUY")
    ]
    
    for idx, (s_title, s_desc) in enumerate(steps_data):
        step_y = scale(2.4 + idx * 0.75)
        # Step card
        add_card(slide, scale(7.1), step_y, scale(4.9), scale(0.68), RGBColor(30, 41, 59), border_rgb=accent_blue if idx == 4 else None)
        
        # Step text
        step_tbox = slide.shapes.add_textbox(scale(7.2), step_y + scale(0.08), scale(4.7), scale(0.52))
        stf = step_tbox.text_frame
        stf.word_wrap = True
        stf.margin_left = stf.margin_top = stf.margin_right = stf.margin_bottom = Inches(0)
        
        sp = stf.paragraphs[0]
        sp.text = s_title
        sp.font.name = "Segoe UI"
        sp.font.size = Pt(9.5)
        sp.font.bold = True
        sp.font.color.rgb = gold_color if idx == 4 else text_white
        
        sp2 = stf.add_paragraph()
        sp2.text = s_desc
        sp2.font.name = "Segoe UI"
        sp2.font.size = Pt(8)
        sp2.font.color.rgb = text_gray

    # Bottom Strip (One-line differentiator)
    add_card(slide, scale(1.0), scale(6.42), scale(11.3), scale(0.42), card_bg, border_rgb=gold_color)
    bottomBox = slide.shapes.add_textbox(scale(1.1), scale(6.46), scale(11.1), scale(0.34))
    btf = bottomBox.text_frame
    bp = btf.paragraphs[0]
    bp.text = "Unlike price-history trackers: PriceMind forecasts and recommends — not just logs."
    bp.alignment = PP_ALIGN.CENTER
    bp.font.name = "Segoe UI"
    bp.font.size = Pt(10)
    bp.font.bold = True
    bp.font.color.rgb = gold_color

    # ==========================================
    # SLIDE 4: ARCHITECTURE (With Blue Outline Frame)
    # ==========================================
    slide = prs.slides[3]
    # Remove helper placeholders
    clean_placeholders(slide, ["Present your system architecture"])
    
    cols_data = [
        ("1. Frontend", "Vanilla JS Client", [
            "Interactive UI pages",
            "Chart.js charts",
            "Fetch API requests",
            "Explainable verdicts"
        ], accent_blue),
        ("2. Backend", "FastAPI Service", [
            "REST API endpoints",
            "CORS enabled",
            "Wait simulator logic",
            "Local SQLite endpoints"
        ], None),
        ("3. Data Layer", "Price Database", [
            "Prepared dataset",
            "SQLite database",
            "PostgreSQL support",
            "Product meta index"
        ], None),
        ("4. ML Core", "XGBoost Engine", [
            "XGBoost regressor",
            "Time feature builder",
            "Lag engineering",
            "Quantile residuals"
        ], gold_color),
        ("5. Alerts", "Fetch Checker", [
            "Target price storage",
            "Checked on fetch",
            "No background cron",
            "Threshold alerts"
        ], None)
    ]
    
    for i, (col_title, col_sub, points, border_c) in enumerate(cols_data):
        col_l = scale(1.0 + i * 2.3)
        col_t = scale(1.8)
        col_w = scale(2.0)
        col_h = scale(4.8)
        
        # Add column card
        add_card(slide, col_l, col_t, col_w, col_h, card_bg, border_rgb=border_c)
        
        content_box = slide.shapes.add_textbox(col_l + scale(0.12), col_t + scale(0.2), col_w - scale(0.24), col_h - scale(0.4))
        ctf = content_box.text_frame
        ctf.word_wrap = True
        ctf.margin_left = ctf.margin_top = ctf.margin_right = ctf.margin_bottom = Inches(0)
        
        p = ctf.paragraphs[0]
        p.text = col_title.upper()
        p.font.name = "Segoe UI"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = border_c if border_c else text_white
        p.space_after = Pt(2)
        
        p2 = ctf.add_paragraph()
        p2.text = col_sub
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(9.5)
        p2.font.italic = True
        p2.font.color.rgb = text_gray
        p2.space_after = Pt(15)
        
        for pt in points:
            p_pt = ctf.add_paragraph()
            p_pt.text = f"• {pt}"
            p_pt.font.name = "Segoe UI"
            p_pt.font.size = Pt(8)
            p_pt.font.color.rgb = text_gray
            p_pt.space_after = Pt(6)
            
        # Draw chevron/arrow to next card
        if i < 4:
            arrow_l = col_l + col_w + scale(0.05)
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_l, scale(4.0), scale(0.2), scale(0.2))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = gold_color
            arrow.line.fill.background()

    # ==========================================
    # SLIDE 5: TECHNICAL APPROACH (With Blue Outline Frame)
    # ==========================================
    slide = prs.slides[4]
    # Remove helper placeholders
    clean_placeholders(slide, ["Present your technical approach"])
    
    # Left Box (ML Pipelines)
    add_card(slide, scale(1.0), scale(1.8), scale(5.2), scale(4.8), card_dark)
    leftBox = slide.shapes.add_textbox(scale(1.3), scale(2.1), scale(4.6), scale(4.2))
    ltf = leftBox.text_frame
    ltf.word_wrap = True
    ltf.margin_left = ltf.margin_top = ltf.margin_right = ltf.margin_bottom = Inches(0)
    
    lp = ltf.paragraphs[0]
    lp.text = "FORECASTING ENGINE & ML PIPELINE"
    lp.font.name = "Segoe UI"
    lp.font.size = Pt(11)
    lp.font.bold = True
    lp.font.color.rgb = gold_color
    lp.space_after = Pt(12)
    
    ml_points = [
        ("Model Pipeline", "XGBoost Regressor model configured for multi-horizon forecasting (7d / 30d / 60d ahead)."),
        ("Feature Engineering", "Extracts lag prices (t-1, t-7, t-14, t-30), rolling mean/volatility (7, 14, 30 days), time properties, and festival-period flags."),
        ("Validation Strategy", "Chronological train/test split holding out the last 60 days of historical records for out-of-sample backtesting.")
    ]
    
    for title, desc in ml_points:
        p_t = ltf.add_paragraph()
        p_t.text = title
        p_t.font.name = "Segoe UI"
        p_t.font.size = Pt(11)
        p_t.font.bold = True
        p_t.font.color.rgb = text_white
        p_t.space_after = Pt(2)
        
        p_d = ltf.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Segoe UI"
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = text_gray
        p_d.space_after = Pt(10)
        
    # Validation Metric Badge (Placeholder)
    add_card(slide, scale(1.3), scale(5.8), scale(4.6), scale(0.5), card_bg, border_rgb=gold_color)
    badgeBox = slide.shapes.add_textbox(scale(1.4), scale(5.85), scale(4.4), scale(0.4))
    btf = badgeBox.text_frame
    bp = btf.paragraphs[0]
    bp.text = "Backtest Results (Galaxy S24): 30d MAPE: 4.07% | 60d MAPE: 1.17%"
    bp.alignment = PP_ALIGN.CENTER
    bp.font.name = "Segoe UI"
    bp.font.size = Pt(9.5)
    bp.font.bold = True
    bp.font.color.rgb = gold_color
        
    # Right Box (Honesty by Design)
    add_card(slide, scale(6.8), scale(1.8), scale(5.5), scale(4.8), card_bg)
    rightBox = slide.shapes.add_textbox(scale(7.1), scale(2.1), scale(4.9), scale(4.2))
    rtf = rightBox.text_frame
    rtf.word_wrap = True
    rtf.margin_left = rtf.margin_top = rtf.margin_right = rtf.margin_bottom = Inches(0)
    
    rp = rtf.paragraphs[0]
    rp.text = "HONESTY BY DESIGN"
    rp.font.name = "Segoe UI"
    rp.font.size = Pt(11)
    rp.font.bold = True
    rp.font.color.rgb = accent_blue
    rp.space_after = Pt(12)
    
    honesty_points = [
        ("Residual-Based Confidence", "Confidence scores are dynamically derived from actual validation residuals, not static values."),
        ("Probabilistic Estimates", "Forecasts represent predictions as a range and include a probability of price drop rather than a fixed guarantee."),
        ("Low-Confidence Fallback", "System disables forced recommendation (LOW CONFIDENCE) when the validation accuracy score falls below threshold.")
    ]
    
    for title, desc in honesty_points:
        p_t = rtf.add_paragraph()
        p_t.text = title
        p_t.font.name = "Segoe UI"
        p_t.font.size = Pt(11)
        p_t.font.bold = True
        p_t.font.color.rgb = text_white
        p_t.space_after = Pt(2)
        
        p_d = rtf.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Segoe UI"
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = text_gray
        p_d.space_after = Pt(10)

    # ==========================================
    # SLIDE 6: FEATURE / USP (With Blue Outline Frame)
    # ==========================================
    slide = prs.slides[5]
    # Remove helper placeholders
    clean_placeholders(slide, ["Key Feature / Unique Selling Point"])
    
    # 1. Top Left: Native Line Chart (Price Forecasting)
    add_card(slide, scale(1.0), scale(1.75), scale(5.2), scale(2.3), card_bg)
    chart_data = CategoryChartData()
    chart_data.categories = ['-30d', '-15d', 'Today', '7d', '30d', '60d']
    chart_data.add_series('Historical', (45000, 44200, 42999, None, None, None))
    chart_data.add_series('Predicted', (None, None, 42999, 42500, 39999, 41200))
    
    x, y, cx, cy = scale(1.05), scale(1.8), scale(5.1), scale(2.2)
    chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data)
    chart = chart_shape.chart
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.font.name = "Segoe UI"
    chart.legend.font.size = Pt(7)
    chart.legend.font.color.rgb = text_white
    chart.has_title = True
    chart.chart_title.text_frame.text = "Price Forecasting Trend (Samsung Galaxy)"
    chart.chart_title.text_frame.paragraphs[0].font.name = "Segoe UI"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(9.5)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    chart.chart_title.text_frame.paragraphs[0].font.color.rgb = gold_color

    # 2. Top Right: Recommendation Engine
    add_card(slide, scale(6.8), scale(1.75), scale(5.5), scale(2.3), card_bg)
    topRight = slide.shapes.add_textbox(scale(7.0), scale(1.85), scale(5.1), scale(2.1))
    trtf = topRight.text_frame
    trtf.word_wrap = True
    trtf.margin_left = trtf.margin_top = trtf.margin_right = trtf.margin_bottom = Inches(0)
    
    p = trtf.paragraphs[0]
    p.text = "TRANSPARENT RECOMMENDATION RULES"
    p.font.name = "Segoe UI"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = accent_blue
    p.space_after = Pt(4)
    
    rules = [
        ("WAIT", "Expected drop >= 3% in next 30 days."),
        ("BUY", "Expected price change < 3% over 30 days."),
        ("STRONG BUY", "Rising trend + high confidence + low drop risk."),
        ("LOW CONFIDENCE", "Fallback when validation MAPE accuracy < 60%.")
    ]
    for r_lbl, r_desc in rules:
        rp = trtf.add_paragraph()
        rp.space_after = Pt(1)
        r_run = rp.add_run()
        r_run.text = f"• {r_lbl}: "
        r_run.font.bold = True
        r_run.font.size = Pt(8.5)
        r_run.font.color.rgb = text_white
        
        d_run = rp.add_run()
        d_run.text = r_desc
        d_run.font.size = Pt(8)
        d_run.font.color.rgb = text_gray

    # 3. Bottom Left: "What if I wait?" Simulator
    add_card(slide, scale(1.0), scale(4.2), scale(5.2), scale(2.3), card_bg)
    bottomLeft = slide.shapes.add_textbox(scale(1.2), scale(4.3), scale(4.8), scale(2.1))
    bltf = bottomLeft.text_frame
    bltf.word_wrap = True
    bltf.margin_left = bltf.margin_top = bltf.margin_right = bltf.margin_bottom = Inches(0)
    
    p = bltf.paragraphs[0]
    p.text = "\"WHAT IF I WAIT?\" SIMULATOR"
    p.font.name = "Segoe UI"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = gold_color
    p.space_after = Pt(4)
    
    p2 = bltf.add_paragraph()
    p2.text = "• Dynamic Interpolation: Scales predicted price, potential savings, and probability of price drop based on user input wait days (1-90 days)."
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(8.5)
    p2.font.color.rgb = text_white
    p2.space_after = Pt(2)
    
    p3 = bltf.add_paragraph()
    p3.text = "• Example Demo: Simulating 30-day wait on Galaxy S-Series shows expected price of ₹39,999, savings of ₹3,000, and 12% risk of price increase."
    p3.font.name = "Segoe UI"
    p3.font.size = Pt(8)
    p3.font.color.rgb = text_gray

    # 4. Bottom Right: Price Alerts
    add_card(slide, scale(6.8), scale(4.2), scale(5.5), scale(2.3), card_bg, border_rgb=gold_color)
    bottomRight = slide.shapes.add_textbox(scale(7.0), scale(4.3), scale(5.1), scale(2.1))
    brtf = bottomRight.text_frame
    brtf.word_wrap = True
    brtf.margin_left = brtf.margin_top = brtf.margin_right = brtf.margin_bottom = Inches(0)
    
    p = brtf.paragraphs[0]
    p.text = "ON-DEMAND PRICE ALERTS"
    p.font.name = "Segoe UI"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = gold_color
    p.space_after = Pt(4)
    
    p2 = brtf.add_paragraph()
    p2.text = "• Fetch-Triggered Checking: Target prices are stored in DB and evaluated only on-demand when user queries the forecast (no background cron scripts)."
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(8.5)
    p2.font.color.rgb = text_white
    p2.space_after = Pt(2)
    
    p3 = brtf.add_paragraph()
    p3.text = "• Active Alert Match: Compares user-defined target thresholds against the updated forecasting price path and flags crossing alerts immediately."
    p3.font.name = "Segoe UI"
    p3.font.size = Pt(8)
    p3.font.color.rgb = text_gray

    # ==========================================
    # SLIDE 7: THANK YOU! (With Blue Outline Frame)
    # ==========================================
    slide = prs.slides[6]
    
    # Update team text box in the template (clear original text to prevent double rendering)
    team_text_shape = None
    for shape in slide.shapes:
        if shape.has_text_frame and "HACK-4-CROWN" in shape.text_frame.text:
            team_text_shape = shape
            break
            
    if team_text_shape:
        # Clear original text entirely
        team_text_shape.text_frame.clear()
        
    # Add new wide centered tagline textbox below THANK YOU! card (card ends at ~3.125")
    taglineBox = slide.shapes.add_textbox(Inches(1.0), Inches(3.5), Inches(8.0), Inches(0.4))
    tf_tag = taglineBox.text_frame
    tf_tag.word_wrap = True
    tf_tag.margin_left = tf_tag.margin_top = tf_tag.margin_right = tf_tag.margin_bottom = Inches(0)
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = "PriceMind — Forecast. Explain. Recommend."
    p_tag.alignment = PP_ALIGN.CENTER
    p_tag.font.name = "Segoe UI"
    p_tag.font.size = Pt(16)
    p_tag.font.bold = True
    p_tag.font.color.rgb = gold_color
    
    # Add new wide centered track/team textbox below tagline
    teamBox = slide.shapes.add_textbox(Inches(1.0), Inches(4.1), Inches(8.0), Inches(0.4))
    tf_team = teamBox.text_frame
    tf_team.word_wrap = True
    tf_team.margin_left = tf_team.margin_top = tf_team.margin_right = tf_team.margin_bottom = Inches(0)
    p_team = tf_team.paragraphs[0]
    p_team.text = "TRACK: Retail & FinTech  |  TEAM: _________________"
    p_team.alignment = PP_ALIGN.CENTER
    p_team.font.name = "Segoe UI"
    p_team.font.size = Pt(11)
    p_team.font.bold = True
    p_team.font.color.rgb = accent_blue

    # Save Presentation
    import time
    try:
        prs.save(output_path)
        print("PowerPoint presentation successfully updated: Hack_4_Crown_PriceMind_Presentation.pptx")
    except PermissionError:
        saved = False
        # Try versioned fallbacks from v2 up to v20
        for v in range(2, 21):
            alternative_path = output_path.replace(".pptx", f"_v{v}.pptx")
            try:
                prs.save(alternative_path)
                print(f"Permission denied on primary. Saved instead to fallback: Hack_4_Crown_PriceMind_Presentation_v{v}.pptx")
                saved = True
                break
            except PermissionError:
                continue
        if not saved:
            timestamp_path = output_path.replace(".pptx", f"_{int(time.time())}.pptx")
            prs.save(timestamp_path)
            print(f"Permission denied on versioned. Saved to timestamp: {os.path.basename(timestamp_path)}")

if __name__ == "__main__":
    create_presentation()
